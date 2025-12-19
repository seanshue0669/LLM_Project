# line_bot/file_extractor.py
import os
from typing import Optional
import PyPDF2
from docx import Document
from pptx import Presentation


class FileExtractor:
    """Extract text content from various file formats"""
    
    SUPPORTED_FORMATS = ['.txt', '.pdf', '.docx', '.pptx']
    MAX_FILE_SIZE = 5 * 1024 * 1024  # 5MB
    
    @staticmethod
    def validate_file(file_path: str) -> tuple[bool, Optional[str]]:
        """
        Validate file format and size
        
        Returns:
            (is_valid, error_message)
        """
        if not os.path.exists(file_path):
            return False, "檔案不存在"
        
        # Check file size
        file_size = os.path.getsize(file_path)
        if file_size > FileExtractor.MAX_FILE_SIZE:
            return False, f"檔案大小超過限制 (最大 5MB，目前 {file_size / 1024 / 1024:.2f}MB)"
        
        # Check file extension
        _, ext = os.path.splitext(file_path)
        if ext.lower() not in FileExtractor.SUPPORTED_FORMATS:
            return False, f"不支援的檔案格式：{ext}（支援格式：{', '.join(FileExtractor.SUPPORTED_FORMATS)}）"
        
        return True, None
    
    @staticmethod
    def extract(file_path: str) -> str:
        """
        Extract text content from file
        
        Args:
            file_path: Path to the file
            
        Returns:
            Extracted text content
            
        Raises:
            Exception: If extraction fails
        """
        # Validate first
        is_valid, error_msg = FileExtractor.validate_file(file_path)
        if not is_valid:
            raise Exception(error_msg)
        
        _, ext = os.path.splitext(file_path)
        ext = ext.lower()
        
        try:
            if ext == '.txt':
                return FileExtractor._extract_txt(file_path)
            elif ext == '.pdf':
                return FileExtractor._extract_pdf(file_path)
            elif ext == '.docx':
                return FileExtractor._extract_docx(file_path)
            elif ext == '.pptx':
                return FileExtractor._extract_pptx(file_path)
            else:
                raise Exception(f"不支援的檔案格式：{ext}")
        except Exception as e:
            raise Exception(f"檔案內容提取失敗：{str(e)}")
    
    @staticmethod
    def _extract_txt(file_path: str) -> str:
        """Extract text from .txt file"""
        encodings = ['utf-8', 'big5', 'gb2312', 'latin-1']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    content = f.read().strip()
                    if content:
                        return content
            except (UnicodeDecodeError, UnicodeError):
                continue
        
        raise Exception("無法解析文字檔編碼")
    
    @staticmethod
    def _extract_pdf(file_path: str) -> str:
        """Extract text from .pdf file"""
        text_parts = []
        
        with open(file_path, 'rb') as f:
            pdf_reader = PyPDF2.PdfReader(f)
            
            if len(pdf_reader.pages) == 0:
                raise Exception("PDF 檔案沒有內容")
            
            for page in pdf_reader.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text)
        
        content = '\n'.join(text_parts).strip()
        
        if not content:
            raise Exception("PDF 檔案無法提取文字（可能是圖片掃描檔）")
        
        return content
    
    @staticmethod
    def _extract_docx(file_path: str) -> str:
        """Extract text from .docx file"""
        doc = Document(file_path)
        
        text_parts = []
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if text:
                text_parts.append(text)
        
        content = '\n'.join(text_parts).strip()
        
        if not content:
            raise Exception("Word 檔案沒有文字內容")
        
        return content
    
    @staticmethod
    def _extract_pptx(file_path: str) -> str:
        """Extract text from .pptx file"""
        prs = Presentation(file_path)
        
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        text_parts.append(text)
        
        content = '\n'.join(text_parts).strip()
        
        if not content:
            raise Exception("PowerPoint 檔案沒有文字內容")
        
        return content